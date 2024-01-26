from pptx import Presentation
import glob
import pandas as pd

###################################################################

# Chargement du powerpoint pour extraction de toutes les informations du documents contenant du texte.

pptx_file = "C:\\Users\\Ilyes\\Desktop\\Données\\Exemple FASEP.pptx"

prs = Presentation(pptx_file) 

processed_ppt = []

print(pptx_file)
print("----------------------")

def extract_text_from_shape(shape, result_list):
    if hasattr(shape, "text"):
        result_list.append(shape.text) 
    elif shape.has_table:
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    result_list.append(paragraph.text)

for slide in prs.slides:
    for shape in slide.shapes:
        extract_text_from_shape(shape, processed_ppt)



##################################################################
        
"""
Processing du texte extrait pour garder les informations les plus importantes 
et une structure globalement normalisée pour éviter des erreurs

"""            
processed_ppt = [word.lower() for word in processed_ppt]

processed_ppt

# Charger la liste de stopwords
stopwords_file = "C:\\Users\\Ilyes\\Desktop\\Données\\stop_words_french.txt"
with open(stopwords_file, "r", encoding="utf-8") as f:
    stopwords_list = f.read().splitlines()

# Prétraitement du texte  
cleaned_text = []
for line in processed_ppt:

    # Supprimer espaces en début et fin
    cleaned_line = line.strip()
    
    # Convertir en minuscules
    cleaned_line = cleaned_line.lower()  

    # Supprimer mots vides
    cleaned_words = [word for word in cleaned_line.split() if word not in stopwords_list]

    # Joindre mots sans espaces entre
    cleaned_line = "_".join(cleaned_words)

    cleaned_text.append(cleaned_line)

# Texte nettoyé  
cleaned_text_combined = "_".join(cleaned_text)



####################################################################
"""
Conversion de la liste de texte génerée, en data frame pandas simple à manipuler pour mettre en place une logique de recherche des termes
et d'impression de la valeur rangée dans la ligne du dessous

"""
df = pd.DataFrame(cleaned_text)
df
def extraire_valeurs(df):
    montant_fasep = None
    date_convention_nataxis = None
    avis_versement_intermediaire = None
    
    for index, row in df.iterrows():
        if 'montant_fasep' in row.values:
            montant_fasep = df.iloc[index + 1].values[0]
        
        if 'date_signature_convention_natixis' in row.values:
            date_convention_nataxis = df.iloc[index + 1].values[0]
        
        if 'avis_versement_intermédiaire' in row.values and index < 55 :
            avis_versement_intermediaire = df.iloc[index + 1].values[0]
    
    return montant_fasep, date_convention_nataxis, avis_versement_intermediaire


montant_fasep, date_convention_nataxis, avis_versement_intermediaire = extraire_valeurs(df)

# Impression des résultats dans la console de sortie

print("Montant FASEP:", montant_fasep)
print("Date convention Nataxis:", date_convention_nataxis)
print("Avis versement intermédiaire:", avis_versement_intermediaire)

########################################################################

